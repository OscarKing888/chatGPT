import os
from pptx import Presentation
import argparse
import re
import shutil
import tempfile
from zipfile import ZipFile
import string

def parse_page_ranges(page_ranges_str):
    """解析页面范围字符串，返回页面索引列表和范围分组"""
    if not page_ranges_str:
        return None, None
    
    try:
        print(f"解析页面范围: {page_ranges_str}")
        page_indices = []
        range_groups = []
        ranges = page_ranges_str.split(',')
        
        for range_str in ranges:
            range_str = range_str.strip()
            print(f"处理范围: {range_str}")
            if '-' in range_str:
                start, end = map(int, range_str.split('-'))
                if start > end:
                    raise ValueError(f"无效的页面范围：起始页({start})大于结束页({end})")
                if start < 1:
                    raise ValueError(f"无效的页面范围：起始页({start})小于1")
                # 转换为0-based索引
                current_range = list(range(start - 1, end))
                page_indices.extend(current_range)
                range_groups.append(current_range)
            else:
                # 单个页面
                page_num = int(range_str)
                if page_num < 1:
                    raise ValueError(f"无效的页面号：{page_num}")
                page_indices.append(page_num - 1)
                range_groups.append([page_num - 1])
        
        # 去重并排序
        page_indices = sorted(list(set(page_indices)))
        print(f"解析到的页面范围组: {len(range_groups)}")
            
        return page_indices, range_groups
    except ValueError as e:
        print(f"错误：{str(e)}")
        return None, None
    except Exception as e:
        print(f"解析页面范围时出错：{str(e)}")
        return None, None

def get_slide_title(slide):
    """获取幻灯片的标题文本"""
    title_text = ""
    
    # 查找标题形状
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            # 如果形状有文本，且文本不为空，我们假定它是标题
            if shape.text.strip():
                title_text = shape.text.strip()
                break
    
    # 如果没找到标题，返回空字符串
    return title_text

def sanitize_filename(filename):
    """清理文件名，将非法字符替换为下划线"""
    # 替换换行符为短横线
    filename = filename.replace('\n', '').replace('\r', '').replace('\x0b', '-')
    
    # 只保留可打印的ASCII字符和中文字符
    result = ""
    for char in filename:
        # 检查是否是可打印的ASCII字符或中文字符
        if char in string.printable or ord(char) > 127:
            result += char
        else:
            result += '_'
    
    # 替换Windows文件名不允许的字符
    invalid_chars = r'[\\/*?:"<>|]'
    result = re.sub(invalid_chars, '_', result)
    
    # 移除开头和结尾的空格和点号
    result = result.strip(" .")
    
    # 限制文件名长度
    if len(result) > 150:
        result = result[:147] + "..."
    
    # 如果文件名为空，返回默认值
    if not result:
        result = "untitled"
    
    # 打印调试信息
    print(f"原始标题: {filename}")
    print(f"处理后的文件名: {result}")
    
    return result

def split_pptx(input_file, output_dir, page_ranges_str=None):
    """将PPTX文件按章节或指定页面范围切分"""
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"输入文件 {input_file} 不存在")
    
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # 加载PPTX文件
    try:
        prs = Presentation(input_file)
        total_slides = len(prs.slides)
        print(f"PPT总页数：{total_slides}")
    except Exception as e:
        raise Exception(f"无法打开PPTX文件：{str(e)}")
    
    # 获取文件名（不含扩展名）
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    
    # 解析页面范围
    page_indices, range_groups = parse_page_ranges(page_ranges_str)
    
    if page_indices is None:
        # 如果没有指定页面范围，尝试按章节切分
        current_section = []
        section_count = 0
        
        for i, slide in enumerate(prs.slides):
            current_section.append(i)
            
            # 检查是否是章节标题（通过检查标题文本是否包含"章"或"节"）
            is_section_title = False
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.lower()
                    if "章" in text or "节" in text:
                        is_section_title = True
                        break
            
            # 如果是章节标题且不是第一页，保存当前章节
            if is_section_title and i > 0:
                extract_slides(input_file, current_section[:-1], output_dir, base_name, section_count, prs)
                current_section = [i]
                section_count += 1
        
        # 保存最后一个章节
        if current_section:
            extract_slides(input_file, current_section, output_dir, base_name, section_count, prs)
    else:
        # 检查页面范围是否有效
        if max(page_indices) >= total_slides:
            raise ValueError(f"指定的页面范围超出PPT总页数({total_slides})")
        
        # 按指定页面范围切分
        print(f"准备切分 {len(range_groups)} 个章节")
        for i, range_group in enumerate(range_groups):
            print(f"切分章节 {i+1}，页面范围: {[idx+1 for idx in range_group]}")
            extract_slides(input_file, range_group, output_dir, base_name, i, prs)

def extract_slides(input_file, slide_indices, output_dir, base_name, section_count, prs=None):
    """提取指定索引的幻灯片并保存为新的PPTX文件"""
    if not slide_indices:
        print(f"章节 {section_count + 1} 没有有效的幻灯片，跳过")
        return
    
    # 如果还没有加载演示文稿，现在加载
    if prs is None:
        prs = Presentation(input_file)
    
    # 获取第一页的标题作为文件名
    first_slide_index = slide_indices[0]
    if first_slide_index < len(prs.slides):
        title_text = get_slide_title(prs.slides[first_slide_index])
        if title_text:
            section_name = sanitize_filename(title_text)
        else:
            section_name = f"section_{section_count + 1}"
    else:
        section_name = f"section_{section_count + 1}"
    
    # 创建一个新的演示文稿
    temp_pptx = os.path.join(output_dir, f"{base_name}_{section_name}.pptx")
    
    try:
        # 复制原始文件
        shutil.copy(input_file, temp_pptx)
        
        # 打开新文件
        new_prs = Presentation(temp_pptx)
        
        # 获取要保留的幻灯片索引（倒序排列以避免删除影响索引）
        keep_indices = sorted(slide_indices)
        delete_indices = [i for i in range(len(new_prs.slides)) if i not in keep_indices]
        delete_indices.sort(reverse=True)
        
        # 删除不需要的幻灯片
        for idx in delete_indices:
            if idx < len(new_prs.slides):
                xml_slides = new_prs.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[idx])
        
        # 保存结果
        new_prs.save(temp_pptx)
        print(f"已保存章节 {section_count + 1} 到: {temp_pptx}")
    except Exception as e:
        print(f"处理文件时出错：{str(e)}")
        # 如果出错，尝试使用序号作为文件名
        try:
            fallback_name = f"{base_name}_section_{section_count + 1}.pptx"
            temp_pptx = os.path.join(output_dir, fallback_name)
            
            # 复制原始文件
            shutil.copy(input_file, temp_pptx)
            
            # 打开新文件
            new_prs = Presentation(temp_pptx)
            
            # 获取要保留的幻灯片索引
            keep_indices = sorted(slide_indices)
            delete_indices = [i for i in range(len(new_prs.slides)) if i not in keep_indices]
            delete_indices.sort(reverse=True)
            
            # 删除不需要的幻灯片
            for idx in delete_indices:
                if idx < len(new_prs.slides):
                    xml_slides = new_prs.slides._sldIdLst
                    slides = list(xml_slides)
                    xml_slides.remove(slides[idx])
            
            # 保存结果
            new_prs.save(temp_pptx)
            print(f"使用备用文件名保存章节 {section_count + 1} 到: {temp_pptx}")
        except Exception as e2:
            raise Exception(f"保存文件时出错: {str(e2)}")

def main():
    parser = argparse.ArgumentParser(description="PPTX文件切分工具")
    parser.add_argument("input_file", help="输入的PPTX文件路径")
    parser.add_argument("output_dir", help="输出目录路径")
    parser.add_argument("--page-ranges", help="页面范围，格式如：1-3,5-8,9-10")
    
    args = parser.parse_args()
    
    try:
        print(f"输入文件: {args.input_file}")
        print(f"输出目录: {args.output_dir}")
        print(f"页面范围: {args.page_ranges}")
        
        split_pptx(args.input_file, args.output_dir, args.page_ranges)
        print("切分完成！")
    except Exception as e:
        print(f"发生错误: {str(e)}")
        exit(1)

if __name__ == "__main__":
    main() 